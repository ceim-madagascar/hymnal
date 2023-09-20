import click
import re
from pptx import Presentation
import glob
from typing import Tuple
from hymnal.utils import slugify
from hymnal.db import get_db


def init_app(app):
    app.cli.add_command(import_command)


PAGE_NUMBER_PATTERN = r'(\d+)(\/\d+)?'


@click.command("import")
@click.argument("path")
def import_command(path):
    files_that_didnt_work = []
    # extract from pptx files

    hymns = {}
    for pptx_file in glob.glob(path + '**/*.pptx', recursive=True):
        try:
            title, content = extract_title_and_content_from_pptx_file(pptx_file)
            hymns[slugify(title)] = (title, content)
        except Exception:
            raise 
            # files_that_didnt_work.append(pptx_file)
            # continue

    # extract from txt files
    for txt_file in glob.glob(path + '**/*.txt', recursive=True):
        try:
            title, content = extract_title_and_content_from_txt_file(txt_file)
            hymns[slugify(title)] = (title, content)
        except Exception:
            raise
            # files_that_didnt_work.append(txt_file)
            # continue

    if files_that_didnt_work:
        raise Exception("The following files could not be processed:" + str(files_that_didnt_work))

    # process hymns
    for slug, (title, content) in hymns.items():
        print("processing " + title)
        db = get_db()
        db.execute(
            'INSERT INTO hymns (title, slug, content) VALUES (?, ?, ?)',
            (title, slug, content)
        )
        db.commit()


def extract_title_and_content_from_pptx_file(filepath) -> Tuple:
    title = filepath.split("/")[-1].strip('.pptx')
    print("processing " + filepath)
    prs = Presentation(filepath)

    content_slides = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if txt := getattr(shape, "text", "").strip():
                if not _line_can_be_skipped(txt):
                    # split the text into lines and remove any line breaks
                    # standardize line break to \n
                    lines = [d.strip() for d in txt.splitlines()]
                    content_slides.append("\n".join(lines))

    content = "\n\n".join(content_slides)
    return title, content


def extract_title_and_content_from_txt_file(filepath):
    title = filepath.split("/")[-1].strip('.txt').strip('.ppt')
    print("processing " + filepath)

    with open(filepath, 'r') as f:
        content = "".join([
            line for line in f.readlines() 
            if not _line_can_be_skipped(line)
        ])
    return title, content


def _line_can_be_skipped(line) -> bool:
    return any([
        line.strip() == "",
        re.match(PAGE_NUMBER_PATTERN, line.strip()),
        "ceim" in line.lower()
    ])
